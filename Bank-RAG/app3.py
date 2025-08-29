"""
Banking AI Chatbot - single-file Streamlit app
Features:
 - Admin (upload/train) and User (ask questions + feedback) roles
 - Uploads saved to bank_docs/
 - Vectorstore persisted to bank_vectorstore/
 - Feedback saved to bank_feedback.csv
 - Uses LangChain + HuggingFaceEmbeddings + FAISS + ChatOpenAI (configurable)
"""

import os
import sys
import types
import time
import json
import httpx
import truststore
import streamlit as st
from datetime import datetime
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation

from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain.prompts import PromptTemplate
from langchain.chains import RetrievalQA
from langchain.memory import ConversationBufferMemory
from langchain_openai import ChatOpenAI

# ----------------------------
# Configuration (edit if needed)
# ----------------------------
DOCS_DIR = "bank_docs"
VECTORSTORE_DIR = "bank_vectorstore"
FEEDBACK_FILE = "bank_feedback.csv"

os.makedirs(DOCS_DIR, exist_ok=True)
os.makedirs(VECTORSTORE_DIR, exist_ok=True)

# Admin credentials (for demo). For production, use env vars or a proper auth provider.
ADMIN_USERNAME = os.getenv("BANK_ADMIN_USER", "bankadmin")
ADMIN_PASSWORD = os.getenv("BANK_ADMIN_PASS", "secure123")

# LLM configuration - configure via environment variables or edit below
LLM_BASE_URL = os.getenv("LLM_BASE_URL", "https://genailab.tcs.in")   # replace if needed
LLM_MODEL = os.getenv("LLM_MODEL", "azure_ai/genailab-maas-DeepSeek-V3-0324")
LLM_API_KEY = os.getenv("LLM_API_KEY", "")  # set in env for production

# Small workaround for Torch + Streamlit watcher issues (if torch is installed)
try:
    import torch
    if not hasattr(torch, "classes"):
        torch.classes = types.SimpleNamespace()
        sys.modules["torch.classes"] = torch.classes
except Exception:
    pass

# SSL fix for corporate environments
try:
    truststore.inject_into_ssl()
except Exception:
    pass

# ----------------------------
# LLM client (ChatOpenAI wrapper)
# ----------------------------
# Use an httpx client that skips cert verification only if necessary (enterprise networks).
_client = httpx.Client(verify=False)

llm = ChatOpenAI(
    base_url=LLM_BASE_URL,
    model=LLM_MODEL,
    api_key=LLM_API_KEY or None,
    http_client=_client,
    temperature=0.2
)

# ----------------------------
# Document processing helpers
# ----------------------------
def save_uploaded_file(uploaded_file):
    """Save uploaded Streamlit file immediately to DOCS_DIR (overwrite if exists)."""
    dest = os.path.join(DOCS_DIR, uploaded_file.name)
    with open(dest, "wb") as f:
        f.write(uploaded_file.getbuffer())
    return dest

def extract_text_from_pdf(path):
    reader = PdfReader(path)
    text_parts = []
    for page in reader.pages:
        txt = page.extract_text()
        if txt:
            text_parts.append(txt)
    return "\n".join(text_parts)

def extract_text_from_docx(path):
    doc = DocxDocument(path)
    return "\n".join([p.text for p in doc.paragraphs])

def extract_text_from_pptx(path):
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
    return "\n".join(texts)

def extract_text(path):
    lower = path.lower()
    if lower.endswith(".pdf"):
        return extract_text_from_pdf(path)
    if lower.endswith(".docx"):
        return extract_text_from_docx(path)
    if lower.endswith(".pptx"):
        return extract_text_from_pptx(path)
    return ""

# ----------------------------
# Chunking & Vectorstore helpers
# ----------------------------
def split_texts_into_chunks(doc_texts, chunk_size=1000, chunk_overlap=200):
    splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap, length_function=len)
    chunks = []
    metadatas = []
    for text, metadata in doc_texts:
        pieces = splitter.split_text(text)
        chunks.extend(pieces)
        metadatas.extend([metadata] * len(pieces))
    return chunks, metadatas

def create_and_persist_vectorstore(chunks, metadatas, persist_dir=VECTORSTORE_DIR):
    embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-mpnet-base-v2")
    vs = FAISS.from_texts(texts=chunks, embedding=embeddings, metadatas=metadatas)
    os.makedirs(persist_dir, exist_ok=True)
    vs.save_local(persist_dir)
    return vs

def load_vectorstore_if_exists(persist_dir=VECTORSTORE_DIR):
    if not os.path.exists(persist_dir) or not os.listdir(persist_dir):
        raise FileNotFoundError("Vectorstore not found")
    embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-mpnet-base-v2")
    vs = FAISS.load_local(persist_dir, embeddings, allow_dangerous_deserialization=True)
    return vs

# ----------------------------
# QA chain helper
# ----------------------------
def get_qa_chain_from_vectorstore(vectorstore):
    template = """
You are a helpful, accurate Banking Assistant. Use only the retrieved document context to answer the user's question.
If the answer is not in the documents, say "I don't know" and recommend contacting the bank support team.

Context: {context}
Question: {question}

Answer concisely with banking-safe guidance.
"""
    prompt = PromptTemplate.from_template(template)
    memory = ConversationBufferMemory(memory_key="chat_history", return_messages=True)
    qa = RetrievalQA.from_chain_type(
        llm=llm,
        chain_type="stuff",
        retriever=vectorstore.as_retriever(),
        memory=memory,
        chain_type_kwargs={"prompt": prompt}
    )
    return qa

# ----------------------------
# Feedback persistence
# ----------------------------
def save_feedback_csv(role_label, question, answer, comment, file=FEEDBACK_FILE):
    header_needed = not os.path.exists(file)
    ts = datetime.utcnow().isoformat()
    row = {
        "timestamp_utc": ts,
        "role": role_label,
        "question": question,
        "answer": answer,
        "comment": comment
    }
    line = ",".join('"' + (str(row[k]).replace('"', '""')) + '"' for k in row.keys())
    if header_needed:
        with open(file, "w", encoding="utf-8") as f:
            f.write(",".join(row.keys()) + "\n")
            f.write(line + "\n")
    else:
        with open(file, "a", encoding="utf-8") as f:
            f.write(line + "\n")

def read_feedback(file=FEEDBACK_FILE):
    if not os.path.exists(file):
        return []
    with open(file, "r", encoding="utf-8") as f:
        lines = f.read().splitlines()
    if not lines:
        return []
    headers = [h.strip() for h in lines[0].split(",")]
    rows = []
    for ln in lines[1:]:
        # naive CSV parse for display (CSV library could do better)
        parts = ln.split('","')
        # remove leading/trailing quotes
        parts = [p.strip().strip('"') for p in parts]
        row = dict(zip(headers, parts))
        rows.append(row)
    return rows

# ----------------------------
# Streamlit UI
# ----------------------------
st.set_page_config(page_title="Banking AI Chatbot", layout="wide")
st.title("🏦 Banking AI Chatbot")
st.write("Ask questions about bank policies, procedures, and FAQs. Admins can upload/update bank documents.")

# Sidebar: role + admin login + sidebar feedback
st.sidebar.header("Role & Admin")
role = st.sidebar.radio("Choose role", ["User", "Admin"])

# Sidebar feedback box (always visible)
st.sidebar.markdown("---")
st.sidebar.subheader("📝 Sidebar Feedback")
sidebar_comment = st.sidebar.text_area("Any additional comments or issues?")
if st.sidebar.button("Submit Sidebar Feedback"):
    save_feedback_csv(role_label=role, question="N/A (sidebar)", answer="N/A", comment=sidebar_comment)
    st.sidebar.success("Thank you — sidebar feedback saved.")

# Admin login flow
if role == "Admin":
    if "admin_logged_in" not in st.session_state:
        st.session_state.admin_logged_in = False

    if not st.session_state.admin_logged_in:
        st.sidebar.subheader("🔐 Admin Login")
        admin_user = st.sidebar.text_input("Username", key="admin_user")
        admin_pass = st.sidebar.text_input("Password", type="password", key="admin_pass")
        if st.sidebar.button("Login as Admin"):
            if admin_user == ADMIN_USERNAME and admin_pass == ADMIN_PASSWORD:
                st.session_state.admin_logged_in = True
                st.sidebar.success("Logged in as admin")
                st.experimental_rerun()
            else:
                st.sidebar.error("Invalid admin credentials")
        st.stop()

    # Admin UI: Upload & Train
    st.header("🔑 Admin Portal — Upload & Manage Bank Documents")
    st.info("Upload bank policy docs (PDF / DOCX / PPTX). Files are saved to the server and can be indexed for QA.")

    uploaded_files = st.file_uploader("Upload bank documents (PDF/DOCX/PPTX)", type=["pdf", "docx", "pptx"], accept_multiple_files=True)
    if uploaded_files:
        saved = []
        for f in uploaded_files:
            saved_path = save_uploaded_file(f)
            saved.append(os.path.basename(saved_path))
        st.success(f"Saved {len(saved)} file(s): {', '.join(saved)}")
        # Auto-index the newly uploaded files immediately
        with st.spinner("Indexing uploaded documents..."):
            docs_for_index = []
            for name in saved:
                path = os.path.join(DOCS_DIR, name)
                text = extract_text(path)
                if text and text.strip():
                    docs_for_index.append((text, {"source": name}))
            if docs_for_index:
                chunks, metas = split_texts_into_chunks(docs_for_index)
                create_and_persist_vectorstore(chunks, metas, persist_dir=VECTORSTORE_DIR)
                # store in session so users can use immediately
                try:
                    st.session_state.vectorstore = load_vectorstore_if_exists(VECTORSTORE_DIR)
                except Exception:
                    st.session_state.vectorstore = None
                st.success("Indexing finished and vectorstore updated.")

    # Show stored documents and allow retrain on selected
    st.subheader("Stored Documents (bank_docs/)")
    stored = sorted([f for f in os.listdir(DOCS_DIR) if f.lower().endswith((".pdf",".docx",".pptx"))])
    if stored:
        st.write("\n".join(f"- {s}" for s in stored))
        to_train = st.multiselect("Select stored documents to (re)index", stored, default=stored)
        if st.button("Train Selected Documents"):
            with st.spinner("Building vectorstore from selected documents..."):
                docs_for_index = []
                for name in to_train:
                    path = os.path.join(DOCS_DIR, name)
                    text = extract_text(path)
                    if text and text.strip():
                        docs_for_index.append((text, {"source": name}))
                if docs_for_index:
                    chunks, metas = split_texts_into_chunks(docs_for_index)
                    create_and_persist_vectorstore(chunks, metas, persist_dir=VECTORSTORE_DIR)
                    st.session_state.vectorstore = load_vectorstore_if_exists(VECTORSTORE_DIR)
                    st.success("Vectorstore rebuilt from selected documents.")
            st.experimental_rerun()
    else:
        st.info("No stored bank documents yet. Upload first.")

    # Admin: view feedback
    st.subheader("User Feedback (bank_feedback.csv)")
    feedback_rows = read_feedback(FEEDBACK_FILE)
    if feedback_rows:
        for i, row in enumerate(feedback_rows[::-1], 1):  # show latest first
            st.markdown(f"**{i}.** [{row.get('timestamp_utc','')}] Role: {row.get('role','')}")
            st.write(f"**Q:** {row.get('question','')}")
            st.write(f"**A:** {row.get('answer','')}")
            st.write(f"**Comment:** {row.get('comment','')}")
            st.divider()
    else:
        st.info("No feedback recorded yet.")

    if st.button("Logout"):
        st.session_state.admin_logged_in = False
        st.success("Logged out")
        st.experimental_rerun()

# ---------------- User role UI ----------------
elif role == "User":
    st.header("🤖 Banking Assistant — Ask a question")
    # Load vectorstore from session or disk
    vectorstore = st.session_state.get("vectorstore")
    if vectorstore is None:
        try:
            vectorstore = load_vectorstore_if_exists(VECTORSTORE_DIR)
            st.session_state.vectorstore = vectorstore
        except Exception:
            vectorstore = None

    if vectorstore is None:
        st.warning("No bank documents indexed yet. Please ask an admin to upload and index bank documents.")
        # still show small uploader for power-users? We'll not allow indexing here.
        st.stop()

    qa = get_qa_chain_from_vectorstore(vectorstore)

    # Question input
    user_question = st.text_input("Ask a banking question (e.g., 'How do I apply for a home loan?')", key="user_question_input")

    # Show previous chat (if any)
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = []

    # If user asked
    if user_question:
        with st.spinner("Retrieving answer from bank documents..."):
            try:
                result = qa.invoke({"query": user_question})
                # RetrievalQA sometimes returns dict-like or string; normalize
                answer = result.get("result") if isinstance(result, dict) and "result" in result else str(result)
            except Exception as e:
                answer = f"Error generating answer: {e}"

        # store in chat history
        st.session_state.chat_history.append(("user", user_question))
        st.session_state.chat_history.append(("assistant", answer))
        st.experimental_rerun()

    # Display chat messages
    for i, (role_tag, msg) in enumerate(st.session_state.chat_history):
        with st.chat_message(role_tag):
            st.write(msg)
            # Show feedback box immediately under assistant's message
            if role_tag == "assistant":
                # unique keys to avoid collisions
                fb_key = f"feedback_{i}"
                comment = st.text_area("Feedback for this answer (optional):", key=fb_key, placeholder="Was this answer helpful? Any suggestions?", height=120)
                submit_key = f"submit_{i}"
                if st.button("Submit Feedback", key=submit_key):
                    save_feedback_csv(role_label="User", question=st.session_state.chat_history[i-1][1] if i>0 else "", answer=msg, comment=comment)
                    st.success("Thanks — your feedback was saved.")

    # if no history yet, show placeholder
    if not st.session_state.chat_history:
        st.info("Type a question in the input box above and press Enter to get started.")

# End of app
